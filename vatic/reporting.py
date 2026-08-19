# -------------------------------------*- vatic -*----------------------------
#                               Open Source Risk Analysis
#
#                             Copyright (c) 2026, eggzec
#                          Contact: https://eggzec.github.io/
#
#                         License: GNU General Public License
#                              Version 3, 29 June 2007
#
# ----------------------------------------------------------------------------
#
#  Author(s)
#      Saud Zahir <m.saud.zahir@gmail.com>
#
#  Date
#      7 May 2026
#
#  Description
#      PDF and PowerPoint export utilities for analyses.
#
# ----------------------------------------------------------------------------

from __future__ import annotations

import tempfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

import numpy as np
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from scipy import stats

from vatic.assumption import Assumption
from vatic.logger import get_logger


LOGGER = get_logger(__name__)


@dataclass(frozen=True)
class ReportImage:
    title: str
    file_path: Path


def _base_layout(title: str) -> go.Layout:
    return go.Layout(
        title=title,
        template="plotly_white",
        paper_bgcolor="#FFFFFF",
        plot_bgcolor="#FFFFFF",
        margin=dict(l=70, r=50, t=70, b=60),
        font=dict(family="Segoe UI, sans-serif", size=12, color="#13138C"),
    )


def _fig_histogram(output: np.ndarray) -> go.Figure:
    bins = min(70, max(20, int(np.sqrt(output.size))))
    fig = go.Figure(layout=_base_layout("Forecast Distribution"))
    fig.add_trace(
        go.Histogram(
            x=output,
            nbinsx=bins,
            marker=dict(color="#2323FF", line=dict(color="#0F0F6B", width=1)),
            opacity=0.9,
        )
    )
    fig.update_xaxes(title_text="Outcome")
    fig.update_yaxes(title_text="Frequency")
    return fig


def _fig_cdf(output: np.ndarray) -> go.Figure:
    sorted_values = np.sort(output)
    cumulative = (
        np.arange(1, sorted_values.size + 1, dtype=float) / sorted_values.size
    )
    fig = go.Figure(layout=_base_layout("Cumulative Distribution (CDF)"))
    fig.add_trace(
        go.Scatter(
            x=sorted_values,
            y=cumulative,
            mode="lines",
            line=dict(color="#2323FF", width=2),
        )
    )
    fig.update_xaxes(title_text="Outcome")
    fig.update_yaxes(title_text="Probability", range=[0, 1])
    return fig


def _fig_exceedance(output: np.ndarray) -> go.Figure:
    sorted_values = np.sort(output)
    exceedance = 1.0 - (
        np.arange(1, sorted_values.size + 1, dtype=float) / sorted_values.size
    )
    fig = go.Figure(layout=_base_layout("Exceedance Curve (P(X > x))"))
    fig.add_trace(
        go.Scatter(
            x=sorted_values,
            y=exceedance,
            mode="lines",
            line=dict(color="#C04AFF", width=2),
        )
    )
    fig.update_xaxes(title_text="Threshold")
    fig.update_yaxes(title_text="Probability", range=[0, 1])
    return fig


def _fig_tornado(
    inputs: dict[str, np.ndarray], output: np.ndarray
) -> go.Figure:
    points: list[tuple[str, float]] = []
    output_std = float(np.std(output))
    if output_std != 0.0:
        for name, series in inputs.items():
            if series.size != output.size:
                continue
            if float(np.std(series)) == 0.0:
                continue
            corr = float(np.corrcoef(series, output)[0, 1])
            if np.isnan(corr):
                continue
            points.append((name, corr))
    points.sort(key=lambda item: abs(item[1]), reverse=True)

    fig = go.Figure(layout=_base_layout("Tornado (Sensitivity)"))
    if not points:
        fig.add_annotation(
            text="No varying inputs available for tornado chart",
            x=0.5,
            y=0.5,
            xref="paper",
            yref="paper",
            showarrow=False,
        )
        fig.update_xaxes(visible=False)
        fig.update_yaxes(visible=False)
        return fig

    names = [name for name, _ in points]
    values = [value for _, value in points]
    colorscale = ["#C04AFF" if value < 0 else "#2323FF" for value in values]
    fig.add_trace(
        go.Bar(
            x=values, y=names, orientation="h", marker=dict(color=colorscale)
        )
    )
    fig.update_layout(showlegend=False)
    fig.update_xaxes(
        title_text="Correlation with forecast", range=[-1.0, 1.0], zeroline=True
    )
    fig.update_yaxes(autorange="reversed")
    return fig


def _fig_var_cvar(output: np.ndarray, confidence: float = 0.95) -> go.Figure:
    alpha = 1.0 - confidence
    var_threshold = float(np.quantile(output, alpha))
    tail = output[output <= var_threshold]
    cvar_value = float(np.mean(tail)) if tail.size else var_threshold
    bins = min(70, max(20, int(np.sqrt(output.size))))
    y_max = float(np.max(np.histogram(output, bins=bins)[0]))
    if y_max <= 0.0:
        y_max = 1.0

    fig = go.Figure(
        layout=_base_layout(f"Tail Risk (VaR/CVaR @ {confidence:.0%})")
    )
    fig.add_trace(
        go.Histogram(
            x=output,
            nbinsx=bins,
            marker=dict(color="#9C9CFF", line=dict(color="#2323FF", width=1)),
            opacity=0.85,
        )
    )
    fig.add_trace(
        go.Scatter(
            x=[var_threshold, var_threshold],
            y=[0.0, y_max],
            mode="lines",
            line=dict(color="#C04AFF", width=2, dash="dash"),
            name=f"VaR {confidence:.0%}",
            hovertemplate=f"VaR {confidence:.0%}: {var_threshold:,.4f}<extra></extra>",
        )
    )
    fig.add_trace(
        go.Scatter(
            x=[cvar_value, cvar_value],
            y=[0.0, y_max],
            mode="lines",
            line=dict(color="#772E9E", width=2, dash="dot"),
            name="CVaR",
            hovertemplate=f"CVaR: {cvar_value:,.4f}<extra></extra>",
        )
    )
    fig.update_layout(margin=dict(l=70, r=50, t=120, b=60))
    fig.add_annotation(
        x=0.01,
        y=1.08,
        xref="paper",
        yref="paper",
        xanchor="left",
        showarrow=False,
        text=f"VaR {confidence:.0%}: {var_threshold:,.4f}",
        font=dict(color="#7E3DFF", size=12),
        bgcolor="rgba(255,255,255,0.85)",
    )
    fig.add_annotation(
        x=0.99,
        y=1.08,
        xref="paper",
        yref="paper",
        xanchor="right",
        showarrow=False,
        text=f"CVaR: {cvar_value:,.4f}",
        font=dict(color="#45228C", size=12),
        bgcolor="rgba(255,255,255,0.85)",
    )
    fig.update_xaxes(title_text="Outcome")
    fig.update_yaxes(title_text="Frequency")
    return fig


def _fig_kde(output: np.ndarray) -> go.Figure:
    fig = go.Figure(layout=_base_layout("Probability Density (KDE)"))
    if output.size < 2 or np.isclose(float(np.std(output)), 0.0):
        fig.add_annotation(
            text="KDE unavailable: output has insufficient variance",
            x=0.5,
            y=0.5,
            xref="paper",
            yref="paper",
            showarrow=False,
        )
        fig.update_xaxes(visible=False)
        fig.update_yaxes(visible=False)
        return fig

    kde = stats.gaussian_kde(output)
    x_grid = np.linspace(float(np.min(output)), float(np.max(output)), 400)
    y_density = kde(x_grid)
    fig.add_trace(
        go.Scatter(
            x=x_grid,
            y=y_density,
            mode="lines",
            line=dict(color="#2323FF", width=2.5),
            fill="tozeroy",
            fillcolor="rgba(35,35,255,0.18)",
        )
    )
    fig.update_xaxes(title_text="Outcome")
    fig.update_yaxes(title_text="Density")
    return fig


def _fig_qq_normal(output: np.ndarray) -> go.Figure:
    fig = go.Figure(layout=_base_layout("Q-Q Plot vs Normal"))
    if output.size < 3:
        fig.add_annotation(
            text="Q-Q plot unavailable: need at least 3 samples",
            x=0.5,
            y=0.5,
            xref="paper",
            yref="paper",
            showarrow=False,
        )
        fig.update_xaxes(visible=False)
        fig.update_yaxes(visible=False)
        return fig

    (theoretical, observed), (slope, intercept, corr) = stats.probplot(
        output, dist="norm"
    )
    fit_line = slope * theoretical + intercept
    fig.add_trace(
        go.Scatter(
            x=theoretical,
            y=observed,
            mode="markers",
            marker=dict(size=6, color="#24AEFF", opacity=0.6),
            name="Sample quantiles",
        )
    )
    fig.add_trace(
        go.Scatter(
            x=theoretical,
            y=fit_line,
            mode="lines",
            line=dict(color="#C04AFF", width=2),
            name=f"Reference line (r={corr:.4f})",
        )
    )
    fig.update_xaxes(title_text="Theoretical quantiles (Normal)")
    fig.update_yaxes(title_text="Observed quantiles")
    return fig


def _fig_pareto(output: np.ndarray) -> go.Figure:
    bins = min(24, max(8, int(np.sqrt(output.size))))
    counts, edges = np.histogram(output, bins=bins)
    order = np.argsort(counts)[::-1]
    counts = counts[order]
    labels = [f"{edges[i]:.2f}..{edges[i + 1]:.2f}" for i in order]
    cumulative = (np.cumsum(counts) / np.sum(counts)) * 100

    fig = go.Figure(layout=_base_layout("Pareto (Binned Outcomes)"))
    fig.add_trace(
        go.Bar(
            x=labels, y=counts, marker=dict(color="#24AEFF"), name="Frequency"
        )
    )
    fig.add_trace(
        go.Scatter(
            x=labels,
            y=cumulative,
            mode="lines+markers",
            marker=dict(color="#C04AFF", size=8),
            line=dict(color="#C04AFF", width=2),
            name="Cumulative %",
            yaxis="y2",
        )
    )
    fig.update_layout(
        yaxis=dict(title="Frequency"),
        yaxis2=dict(
            title="Cumulative %",
            overlaying="y",
            side="right",
            range=[0, 105],
            showgrid=False,
        ),
        xaxis=dict(tickangle=45),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, x=0),
    )
    return fig


def _fig_trend(output: np.ndarray) -> go.Figure:
    max_points = 6000
    if output.size > max_points:
        idx = np.linspace(0, output.size - 1, num=max_points, dtype=int)
        sampled = output[idx]
    else:
        sampled = output

    x = np.arange(1, sampled.size + 1)
    running_mean = np.cumsum(sampled) / x
    fig = go.Figure(layout=_base_layout("Trend (Simulation Path)"))
    fig.add_trace(
        go.Scatter(
            x=x,
            y=running_mean,
            mode="lines",
            line=dict(color="#2323FF", width=2.2),
            name="Running mean",
        )
    )
    fig.add_trace(
        go.Scatter(
            x=x,
            y=sampled,
            mode="lines",
            line=dict(color="#9C9CFF", width=1),
            opacity=0.6,
            name="Sample outcome",
        )
    )
    fig.update_xaxes(title_text="Iteration")
    fig.update_yaxes(title_text="Value")
    return fig


def _fig_scatter_primary(
    output: np.ndarray, inputs: dict[str, np.ndarray]
) -> go.Figure:
    fig = go.Figure(layout=_base_layout("Scatter (Primary Input vs Forecast)"))
    if not inputs:
        fig.add_annotation(
            text="Scatter unavailable: no input samples",
            x=0.5,
            y=0.5,
            xref="paper",
            yref="paper",
            showarrow=False,
        )
        fig.update_xaxes(visible=False)
        fig.update_yaxes(visible=False)
        return fig

    best_name = ""
    best_series: np.ndarray | None = None
    best_corr = -1.0
    for name, series in inputs.items():
        if series.size != output.size:
            continue
        if np.isclose(float(np.std(series)), 0.0):
            continue
        corr = float(np.corrcoef(series, output)[0, 1])
        if np.isnan(corr):
            continue
        if abs(corr) > best_corr:
            best_corr = abs(corr)
            best_name = name
            best_series = series

    if best_series is None:
        fig.add_annotation(
            text="Scatter unavailable: no varying input",
            x=0.5,
            y=0.5,
            xref="paper",
            yref="paper",
            showarrow=False,
        )
        fig.update_xaxes(visible=False)
        fig.update_yaxes(visible=False)
        return fig

    max_points = 6000
    if best_series.size > max_points:
        idx = np.linspace(0, best_series.size - 1, num=max_points, dtype=int)
        x = best_series[idx]
        y = output[idx]
    else:
        x = best_series
        y = output

    fig.add_trace(
        go.Scattergl(
            x=x,
            y=y,
            mode="markers",
            marker=dict(color="#4E269E", size=6, opacity=0.45),
            name=best_name,
        )
    )
    fig.update_xaxes(title_text=best_name)
    fig.update_yaxes(title_text="Forecast")
    return fig


def _fig_box(output: np.ndarray, inputs: dict[str, np.ndarray]) -> go.Figure:
    fig = go.Figure(layout=_base_layout("Box Plot (Inputs + Forecast)"))
    groups = dict(inputs)
    groups["Forecast"] = output
    for name, values in groups.items():
        fig.add_trace(
            go.Box(
                x=values,
                name=name,
                boxpoints=False,
                marker=dict(color="#24AEFF"),
                line=dict(color="#2323FF"),
            )
        )
    fig.update_xaxes(title_text="Value")
    return fig


def _fig_violin(output: np.ndarray, inputs: dict[str, np.ndarray]) -> go.Figure:
    fig = go.Figure(layout=_base_layout("Violin Plot (Inputs + Forecast)"))
    groups = dict(inputs)
    groups["Forecast"] = output
    for name, values in groups.items():
        fig.add_trace(
            go.Violin(
                x=values,
                name=name,
                box_visible=True,
                meanline_visible=True,
                points=False,
                line_color="#2323FF",
                fillcolor="rgba(126,61,255,0.30)",
            )
        )
    fig.update_xaxes(title_text="Value")
    return fig


def _fig_rich_statistics(stats: dict[str, float]) -> go.Figure:
    labels = [
        "Samples",
        "Mean",
        "Std Dev",
        "Min / Max",
        "P01 / P05",
        "P50 / P95 / P99",
        "Probability(Output < 0)",
    ]
    values = [
        f"{int(stats['samples']):,}",
        f"{stats['mean']:,.6f}",
        f"{stats['std']:,.6f}",
        f"{stats['min']:,.6f} / {stats['max']:,.6f}",
        f"{stats['p01']:,.6f} / {stats['p05']:,.6f}",
        f"{stats['p50']:,.6f} / {stats['p95']:,.6f} / {stats['p99']:,.6f}",
        f"{stats['prob_loss']:.2%}",
    ]
    fig = go.Figure(
        data=[
            go.Table(
                header=dict(
                    values=["Metric", "Value"],
                    fill_color="#2323FF",
                    font=dict(color="white", size=13),
                    align="left",
                ),
                cells=dict(
                    values=[labels, values],
                    fill_color=["#FFFFFF", "#FFFFFF"],
                    align="left",
                    font=dict(color="#13138C", size=12),
                    height=30,
                ),
            )
        ]
    )
    fig.update_layout(
        title="Rich Statistics",
        template="plotly_white",
        margin=dict(l=40, r=40, t=70, b=20),
        paper_bgcolor="#FFFFFF",
        font=dict(family="Segoe UI, sans-serif", size=12, color="#13138C"),
    )
    return fig


def _render_report_images(
    temp_dir: Path,
    output: np.ndarray,
    inputs: dict[str, np.ndarray],
    stats: dict[str, float],
) -> list[ReportImage]:
    figures: list[tuple[str, go.Figure]] = [
        ("Forecast Distribution", _fig_histogram(output)),
        ("Cumulative Distribution (CDF)", _fig_cdf(output)),
        ("Exceedance Curve", _fig_exceedance(output)),
        ("Tail Risk (VaR/CVaR)", _fig_var_cvar(output)),
        ("Probability Density (KDE)", _fig_kde(output)),
        ("Q-Q Plot vs Normal", _fig_qq_normal(output)),
        ("Pareto (Binned Outcomes)", _fig_pareto(output)),
        ("Trend (Simulation Path)", _fig_trend(output)),
        (
            "Scatter (Primary Input vs Forecast)",
            _fig_scatter_primary(output, inputs),
        ),
        ("Tornado Sensitivity", _fig_tornado(inputs, output)),
        ("Box Plot (Inputs + Forecast)", _fig_box(output, inputs)),
        ("Violin Plot (Inputs + Forecast)", _fig_violin(output, inputs)),
        ("Rich Statistics", _fig_rich_statistics(stats)),
    ]

    images: list[ReportImage] = []
    for idx, (title, fig) in enumerate(figures, start=1):
        image_path = temp_dir / f"chart_{idx:02d}.png"
        fig.write_image(
            str(image_path), format="png", width=1280, height=720, scale=1
        )
        images.append(ReportImage(title=title, file_path=image_path))
    return images


def export_pdf_report(
    destination: Path,
    analysis_name: str,
    formula: str,
    iterations: int,
    assumptions: list[Assumption],
    stats: dict[str, float],
    output: np.ndarray,
    inputs: dict[str, np.ndarray],
    active_forecast_name: str,
    capability_metrics: dict[str, float],
) -> None:
    LOGGER.debug("Exporting PDF report | path=%s", destination)
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(destination), pagesize=A4)
    content = []

    content.append(Paragraph("Vatic Simulation Report", styles["Title"]))
    content.append(Spacer(1, 12))
    content.append(Paragraph(f"Analysis: {analysis_name}", styles["Normal"]))
    content.append(
        Paragraph(
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            styles["Normal"],
        )
    )
    content.append(Paragraph(f"Formula: {formula}", styles["Normal"]))
    content.append(Paragraph(f"Iterations: {iterations:,}", styles["Normal"]))
    content.append(Spacer(1, 12))

    stats_rows = [
        ["Metric", "Value"],
        ["Samples", f"{int(stats['samples']):,}"],
        ["Mean", f"{stats['mean']:,.6f}"],
        ["Std Dev", f"{stats['std']:,.6f}"],
        ["Min / Max", f"{stats['min']:,.6f} / {stats['max']:,.6f}"],
        ["P01 / P05", f"{stats['p01']:,.6f} / {stats['p05']:,.6f}"],
        [
            "P50 / P95 / P99",
            f"{stats['p50']:,.6f} / {stats['p95']:,.6f} / {stats['p99']:,.6f}",
        ],
        ["Probability(Output < 0)", f"{stats['prob_loss']:.2%}"],
    ]
    stats_table = Table(stats_rows, colWidths=[2.3 * inch, 3.7 * inch])
    stats_table.setStyle(
        TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2323FF")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#D3D3FF")),
            (
                "ROWBACKGROUNDS",
                (0, 1),
                (-1, -1),
                [colors.white, colors.HexColor("#FAFAFF")],
            ),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("ALIGN", (0, 0), (-1, -1), "LEFT"),
        ])
    )
    content.append(stats_table)
    content.append(Spacer(1, 12))

    if capability_metrics:
        capability_rows = [["Metric", "Value"]]
        ordered = [
            "lsl",
            "usl",
            "target",
            "Cp",
            "Cpk-lower",
            "Cpk-upper",
            "Cpk",
            "Cpm",
            "Pp",
            "Ppk-lower",
            "Ppk-upper",
            "Ppk",
            "Ppm",
            "Z-LSL",
            "Z-USL",
            "Zst",
            "Zlt",
            "p(N/C)-below",
            "p(N/C)-above",
            "p(N/C)-total",
            "PPM-below",
            "PPM-above",
            "PPM-total",
        ]
        for key in ordered:
            if key not in capability_metrics:
                continue
            value = capability_metrics[key]
            if key.startswith("p(N/C)"):
                text = f"{value:.4%}"
            elif key.startswith("PPM"):
                text = f"{value:,.2f}"
            else:
                text = f"{value:,.6f}"
            capability_rows.append([key, text])

        capability_table = Table(
            capability_rows, colWidths=[2.3 * inch, 3.7 * inch]
        )
        capability_table.setStyle(
            TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2323FF")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#D3D3FF")),
                (
                    "ROWBACKGROUNDS",
                    (0, 1),
                    (-1, -1),
                    [colors.white, colors.HexColor("#FAFAFF")],
                ),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
            ])
        )
        content.append(
            Paragraph(
                f"Capability Metrics ({active_forecast_name})",
                styles["Heading2"],
            )
        )
        content.append(capability_table)
        content.append(Spacer(1, 12))

    assumption_rows = [["Name", "Distribution", "Parameters"]]
    for item in assumptions:
        params = ", ".join(f"{k}={v}" for k, v in item.parameters.items())
        assumption_rows.append([item.name, item.distribution, params])

    assumption_table = Table(
        assumption_rows, colWidths=[1.5 * inch, 1.4 * inch, 3.1 * inch]
    )
    assumption_table.setStyle(
        TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2323FF")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#D3D3FF")),
            (
                "ROWBACKGROUNDS",
                (0, 1),
                (-1, -1),
                [colors.white, colors.HexColor("#FAFAFF")],
            ),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("ALIGN", (0, 0), (-1, -1), "LEFT"),
        ])
    )
    content.append(Paragraph("Assumptions", styles["Heading2"]))
    content.append(assumption_table)

    with tempfile.TemporaryDirectory(prefix="vatic-report-") as tmp:
        images = _render_report_images(Path(tmp), output, inputs, stats)
        for image in images:
            content.append(PageBreak())
            content.append(Paragraph(image.title, styles["Heading2"]))
            content.append(Spacer(1, 8))
            content.append(
                Image(str(image.file_path), width=6.7 * inch, height=3.9 * inch)
            )
        doc.build(content)

    LOGGER.debug("PDF report exported | path=%s", destination)


def export_pptx_report(
    destination: Path,
    analysis_name: str,
    formula: str,
    iterations: int,
    assumptions: list[Assumption],
    stats: dict[str, float],
    output: np.ndarray,
    inputs: dict[str, np.ndarray],
    active_forecast_name: str,
    capability_metrics: dict[str, float],
) -> None:
    LOGGER.debug("Exporting PowerPoint report | path=%s", destination)
    presentation = Presentation()

    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = "Vatic Simulation Report"
    subtitle = cover.placeholders[1]
    subtitle.text = (
        f"Analysis: {analysis_name}\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"Formula: {formula}\n"
        f"Iterations: {iterations:,}"
    )

    summary = presentation.slides.add_slide(presentation.slide_layouts[1])
    summary.shapes.title.text = "Risk Summary"
    body = summary.shapes.placeholders[1].text_frame
    body.clear()
    lines = [
        f"Samples: {int(stats['samples']):,}",
        f"Mean: {stats['mean']:,.6f}",
        f"Std Dev: {stats['std']:,.6f}",
        f"Min / Max: {stats['min']:,.6f} / {stats['max']:,.6f}",
        f"P01 / P05: {stats['p01']:,.6f} / {stats['p05']:,.6f}",
        f"P50 / P95 / P99: {stats['p50']:,.6f} / {stats['p95']:,.6f} / {stats['p99']:,.6f}",
        f"Probability(Output < 0): {stats['prob_loss']:.2%}",
    ]
    for idx, line in enumerate(lines):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(18)

    if capability_metrics:
        capability_slide = presentation.slides.add_slide(
            presentation.slide_layouts[1]
        )
        capability_slide.shapes.title.text = (
            f"Capability Metrics ({active_forecast_name})"
        )
        capability_body = capability_slide.shapes.placeholders[1].text_frame
        capability_body.clear()
        ordered = [
            "lsl",
            "usl",
            "target",
            "Cp",
            "Cpk-lower",
            "Cpk-upper",
            "Cpk",
            "Cpm",
            "Pp",
            "Ppk-lower",
            "Ppk-upper",
            "Ppk",
            "Ppm",
            "Z-LSL",
            "Z-USL",
            "Zst",
            "Zlt",
            "p(N/C)-below",
            "p(N/C)-above",
            "p(N/C)-total",
            "PPM-below",
            "PPM-above",
            "PPM-total",
        ]
        rows: list[str] = []
        for key in ordered:
            if key not in capability_metrics:
                continue
            value = capability_metrics[key]
            if key.startswith("p(N/C)"):
                text = f"{value:.4%}"
            elif key.startswith("PPM"):
                text = f"{value:,.2f}"
            else:
                text = f"{value:,.6f}"
            rows.append(f"{key}: {text}")

        for idx, row in enumerate(rows):
            paragraph = (
                capability_body.paragraphs[0]
                if idx == 0
                else capability_body.add_paragraph()
            )
            paragraph.text = row
            paragraph.level = 0
            paragraph.font.size = Pt(14)

    assumptions_slide = presentation.slides.add_slide(
        presentation.slide_layouts[1]
    )
    assumptions_slide.shapes.title.text = "Assumptions"
    assumptions_body = assumptions_slide.shapes.placeholders[1].text_frame
    assumptions_body.clear()
    for idx, item in enumerate(assumptions):
        params = ", ".join(f"{k}={v}" for k, v in item.parameters.items())
        paragraph = (
            assumptions_body.paragraphs[0]
            if idx == 0
            else assumptions_body.add_paragraph()
        )
        paragraph.text = f"{item.name} | {item.distribution} | {params}"
        paragraph.level = 0
        paragraph.font.size = Pt(14)

    with tempfile.TemporaryDirectory(prefix="vatic-report-") as tmp:
        images = _render_report_images(Path(tmp), output, inputs, stats)
        for image in images:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = image.title

            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            chart_top = Inches(1.0)
            chart_bottom_margin = Inches(0.35)
            side_margin = Inches(0.35)
            avail_w = slide_width - (2 * side_margin)
            avail_h = slide_height - chart_top - chart_bottom_margin

            target_w = avail_w
            target_h = int(target_w * 9 / 16)
            if target_h > avail_h:
                target_h = avail_h
                target_w = int(target_h * 16 / 9)

            left = int((slide_width - target_w) / 2)
            top = int(chart_top + (avail_h - target_h) / 2)
            slide.shapes.add_picture(
                str(image.file_path), left, top, width=target_w, height=target_h
            )

    presentation.save(str(destination))
    LOGGER.debug("PowerPoint report exported | path=%s", destination)
